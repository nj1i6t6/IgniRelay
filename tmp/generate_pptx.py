"""
Generate two PPTX presentations for IgniRelay project:
  PPTX #1: App-focused (current features + future improvements)
  PPTX #2: Overall system framework (macro view)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color Palette ──
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_ORANGE = RGBColor(0xFF, 0x6B, 0x35)
ACCENT_BLUE = RGBColor(0x00, 0x9E, 0xFF)
ACCENT_GREEN = RGBColor(0x00, 0xE6, 0x76)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
DARK_GRAY = RGBColor(0x33, 0x33, 0x3E)
SECTION_BG = RGBColor(0x16, 0x21, 0x3E)
CARD_BG = RGBColor(0x22, 0x2B, 0x45)
WARN_RED = RGBColor(0xFF, 0x45, 0x45)

OUTPUT_DIR = r"C:\Users\radio\Downloads\IDE\CoReM\tmp"

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=14,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft JhengHei'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=13,
                    color=WHITE, bullet_color=ACCENT_ORANGE, font_name='Microsoft JhengHei'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"▸ {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(6)
    return txBox

def add_card(slide, left, top, width, height, title, body_items, title_color=ACCENT_ORANGE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.fill.background()
    shape.shadow.inherit = False

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(8)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.color.rgb = title_color
    p.font.bold = True
    p.font.name = 'Microsoft JhengHei'

    for item in body_items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(11)
        p.font.color.rgb = LIGHT_GRAY
        p.font.name = 'Microsoft JhengHei'
        p.space_before = Pt(3)

def add_table_slide(slide, left, top, width, rows_data, col_widths=None, font_size=10):
    """Add a table to a slide. rows_data[0] is header row."""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(n_rows, n_cols,
                                          Inches(left), Inches(top),
                                          Inches(width), Inches(0.3 * n_rows))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for r, row in enumerate(rows_data):
        for c, cell_text in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = 'Microsoft JhengHei'
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                else:
                    paragraph.font.color.rgb = LIGHT_GRAY
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = ACCENT_BLUE
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CARD_BG if r % 2 == 1 else DARK_GRAY

    return table_shape

# ══════════════════════════════════════════════════════════════
#  PPTX #1: APP 功能與技術 — 烽傳 IgniRelay App
# ══════════════════════════════════════════════════════════════
def generate_app_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── Slide 1: Cover ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, DARK_BG)
    add_textbox(slide, 1, 0.5, 11, 1.2, "烽傳 IgniRelay", 48, ACCENT_ORANGE, True, PP_ALIGN.CENTER)
    add_textbox(slide, 1, 1.7, 11, 0.8, "去中心化災難通訊 — App 功能與技術架構", 28, WHITE, False, PP_ALIGN.CENTER)
    add_textbox(slide, 1, 3.0, 11, 0.5, "「烽火台是古代第一套 Mesh 網路。烽傳，就是把這個概念裝進每個人的口袋。」", 16, LIGHT_GRAY, False, PP_ALIGN.CENTER)
    add_textbox(slide, 1, 5.0, 11, 0.5, "Phase 0 — 概念驗證與 MVP 開發中  |  2026 年 3 月", 14, LIGHT_GRAY, False, PP_ALIGN.CENTER)
    add_textbox(slide, 1, 6.0, 11, 0.5, "Flutter + Dart  |  BLE Mesh  |  Protobuf  |  Ed25519", 13, ACCENT_BLUE, False, PP_ALIGN.CENTER)

    # ── Slide 2: App 定位 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_textbox(slide, 0.5, 0.3, 12, 0.8, "App 定位：災區最後一哩路的手機端入口", 32, ACCENT_ORANGE, True)
    add_textbox(slide, 0.5, 1.2, 12, 0.5, "在 4G/5G、光纖、WiFi 全面癱瘓時，民眾手機透過 BLE Mesh 互聯，形成自組織救災網路", 16, LIGHT_GRAY)

    add_card(slide, 0.5, 2.2, 3.8, 2.5, "🎯 核心目標",
             ["離線通訊：完全不依賴行動網路",
              "SOS 求救：一鍵廣播至周圍節點",
              "物資媒合：去中心化供需配對",
              "零門檻：安裝 App 即可接入"])

    add_card(slide, 4.6, 2.2, 3.8, 2.5, "📱 技術棧",
             ["Framework: Flutter 3.x / Dart 3.2+",
              "BLE: flutter_blue_plus + Nordic BLE (Android)",
              "序列化: Protobuf 3 二進位編碼",
              "加密: Ed25519 簽章驗證",
              "地圖: flutter_map + MBTiles 離線向量圖",
              "儲存: SQLite + SharedPreferences"])

    add_card(slide, 8.9, 2.2, 3.8, 2.5, "📊 目前狀態 (v0.1.15)",
             ["APK 約 260 MB (含全台離線地圖)",
              "跨廠牌 BLE 互通已驗證",
              "Android: Nordic BLE Library",
              "iOS: Core Bluetooth",
              "已有完整 debug log 系統"])

    add_textbox(slide, 0.5, 5.2, 12, 1.5,
                "Tier 分級：Tier 1 全速 (≥40%) → Tier 2 低頻中繼 (20-40%) → Tier 3 極限求生 (<20%)\n"
                "遲滯設計防止頻繁切換：跌破 40% → Tier 2，需充至 ≥60% 才回 Tier 1",
                13, LIGHT_GRAY)

    # ── Slide 3: BLE Mesh 架構 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_textbox(slide, 0.5, 0.3, 12, 0.8, "BLE Mesh 通訊架構", 32, ACCENT_ORANGE, True)

    add_card(slide, 0.5, 1.3, 6.0, 2.8, "📡 GATT Service 設計 (ble_manager.dart · 774 行)",
             ["主服務 UUID: a4d11949... (UUIDv5 ignirelay.com)",
              "Event Char: 事件傳輸通道",
              "Bloom Char: 同步握手通道",
              "Handshake Char: 實體交接通道",
              "",
              "MTU 協商: 手機→手機 517 bytes",
              "          手機→nRF54H20 ~247 bytes",
              "連線逾時: 10 秒 | 冷卻時間: 60 秒"])

    add_card(slide, 6.8, 1.3, 6.0, 2.8, "🔄 同步 Pipeline",
             ["1. 掃描 → 發現 IgniRelay Service UUID 的節點",
              "2. 連線 → MTU 協商 → 服務發現",
              "3. 讀取對端 Bloom Filter (事件 ID 列表)",
              "4. 比對差集 → 僅傳對方缺少的事件",
              "5. 推送 TriageQueue 中高優先事件",
              "6. 從 DB 補充最近 24h 事件",
              "7. 訂閱 Notify → 接收對端推送"])

    add_card(slide, 0.5, 4.4, 4.0, 2.5, "🛡️ 平台感知路由",
             ["Android → Nordic BLE Library",
              "  (解決 MediaTek/Qualcomm 相容性)",
              "  透過 MethodChannel 呼叫 Kotlin",
              "",
              "iOS → flutter_blue_plus",
              "  (Core Bluetooth 原生穩定)"])

    add_card(slide, 4.8, 4.4, 4.0, 2.5, "🔍 去重機制",
             ["記憶體: Set<String> _seenEvents",
              "Bloom Filter: 事件 ID 列表交換",
              "DB: UNIQUE constraint 防重複",
              "",
              "Bug 5 Fix: 取消標記防孤兒 sync",
              "Bug 6 Fix: OPPO blind relay 模式",
              "Bug 7 Fix: GATT Server Notify 反推"])

    add_card(slide, 9.1, 4.4, 3.7, 2.5, "⚡ 掃描參數",
             ["掃描持續: 30 秒",
              "掃描間隔: 5 秒",
              "最大同時連線: 8",
              "序列化連線 (一台連完再連下一台)",
              "30 秒 timeout + 取消標記"])

    # ── Slide 4: Protobuf 協議 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_textbox(slide, 0.5, 0.3, 12, 0.8, "Protobuf 封包協議 (mesh_protocol.proto)", 32, ACCENT_ORANGE, True)

    table_data = [
        ["EventType", "值", "說明", "路由邊界"],
        ["RESOURCE_REGISTER", "0", "物資登記", "里"],
        ["REQUEST_BROADCAST", "1", "需求廣播 (含 QoS)", "里"],
        ["MATCH_INTENT", "2", "意向匹配", "里"],
        ["PHYSICAL_HANDSHAKE", "3", "物理核銷交割", "端對端"],
        ["HAZARD_MARKER", "4", "動態危險圖層標記", "鄉鎮市區"],
        ["QUARANTINE_VOTE", "5", "惡意節點檢舉投票", "鄉鎮市區"],
        ["MATCH_CANCEL", "6", "釋放配對", "里"],
        ["FIRE_ALARM_RF", "7", "433MHz 住警器火警", "鄉鎮市區"],
    ]
    add_table_slide(slide, 0.5, 1.4, 7.0, table_data, col_widths=[2.5, 0.5, 2.5, 1.5])

    add_card(slide, 8.0, 1.4, 4.8, 3.0, "📦 MeshEvent 結構 (54 行 proto)",
             ["event_id: UUID / SHA-256",
              "sender_pub_key: Ed25519 32 bytes",
              "identity_level: 信任等級 0-3",
              "type + urgency: 事件類型 + 優先級",
              "hlc_timestamp + hlc_counter: HLC",
              "ttl: 剩餘跳數 (預設 10)",
              "origin_lat/lng: 創建者座標 (不可改)",
              "received_lat/lng: 每跳覆寫",
              "payload: 序列化業務資料",
              "signature: Ed25519 簽章 64 bytes"])

    add_card(slide, 8.0, 4.7, 4.8, 2.2, "📋 Sub-Messages",
             ["ResourceData: 物資 ID/類型/數量/位置",
              "RequestData: 需求類型/數量/緊急度",
              "MatchIntentData: 媒合配對雙方公鑰",
              "HazardData: 危險類型/嚴重度/半徑",
              "MedicalSummary: 醫療卡摘要",
              "FireAlarmRfData: RF 頻率/RSSI/座標"])

    # ── Slide 5: CRDT + HLC ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_textbox(slide, 0.5, 0.3, 12, 0.8, "CRDT 離線同步 + HLC 混合邏輯時鐘", 32, ACCENT_ORANGE, True)

    add_card(slide, 0.5, 1.3, 5.8, 3.0, "⏰ HLC 混合邏輯時鐘 (hlc.dart · 89 行)",
             ["三元組: (wallclock, counter, nodeId)",
              "nodeId = 公鑰前 8 bytes hex",
              "",
              "now(): 取得當前 HLC 並推進計數器",
              "merge(remote): 交會強制校時",
              "  → 取 max(本地時間, 本地HLC, 遠端HLC)",
              "  → 防止斷電重置導致時間戳倒退",
              "",
              "compareTo(): timestamp → counter → nodeId"])

    add_card(slide, 6.6, 1.3, 6.2, 3.0, "⚖️ 衝突解決器 (conflict_resolver.dart · 40 行)",
             ["Double-Spending 排解邏輯:",
              "  1. HLC timestamp 較小者勝 (先發先得)",
              "  2. HLC counter 較小者勝",
              "  3. Urgency 高者勝 (SOS_RED > YELLOW > RESOURCE)",
              "  4. Tiebreaker: PubKey 字典序",
              "",
              "設計理念: Op-based CmRDT",
              "  → 事件為 append-only 不可變日誌",
              "  → 不需要 Automerge 等通用 CRDT 庫"])

    add_card(slide, 0.5, 4.6, 12.3, 2.5, "🌐 Zone-Based Geo-Fencing 路由 (mesh_router.dart · 113 行)",
             ["路由策略: INFO/RESOURCE → 里邊界 | SOS_YELLOW/HAZARD → 鄉鎮市區 | SOS_RED 已驗證+identity≥1 → 無邊界限制",
              "特殊豁免: Tier 0 硬體騾子、Android Foreground Data Mule 永遠轉發",
              "Fallback: 離島/資料缺漏 → 距離衰減 (SOS_RED 5倍放寬、YELLOW 5倍、RESOURCE 2倍)",
              "行政區查詢: VillageGeofence.isSameVillageZone() / isSameTownshipZone() + 300m 緩衝",
              "Quarantine: 加權投票 (L0=0.2 / L1=0.5 / L2=0.8 / L3=1.0)，累計 > 3.0 黑名單"])

    # ── Slide 6: 已實作功能 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_textbox(slide, 0.5, 0.3, 12, 0.8, "App 已實作功能總覽", 32, ACCENT_ORANGE, True)

    add_card(slide, 0.5, 1.3, 4.0, 2.5, "🔐 身份與加密",
             ["Ed25519 金鑰對自動生成",
              "持久化儲存 (SharedPreferences)",
              "Trust Ladder: L0 匿名 → L3 實名",
              "Quarantine 投票權重依等級",
              "Payload 簽章驗證"])

    add_card(slide, 4.8, 1.3, 4.0, 2.5, "🗺️ 離線地圖",
             ["flutter_map + MBTiles 向量圖",
              "內建全台 OpenMapTiles 離線地圖",
              "長按標記危險區域 (HazardData)",
              "POI 查詢功能 (poi_query.dart)",
              "MBTiles 載入器 (mbtiles_loader.dart)"])

    add_card(slide, 9.1, 1.3, 3.7, 2.5, "🩺 醫療卡",
             ["離線儲存: 姓名/年齡/血型",
              "過敏原 + 反應描述",
              "目前用藥列表",
              "緊急聯絡人",
              "器官捐贈意願",
              "MedicalSummary 附於 SOS"])

    add_card(slide, 0.5, 4.1, 4.0, 2.5, "📦 物資媒合",
             ["ResourceRegister: 登記物資",
              "RequestBroadcast: 廣播需求",
              "MatchIntent: 雙方意向配對",
              "PhysicalHandshake: 4-PIN 交割",
              "MatchCancel: PIN 失敗取消"])

    add_card(slide, 4.8, 4.1, 4.0, 2.5, "📊 Debug 儀表板",
             ["掃描週期計數",
              "節點發現/連線/同步統計",
              "Bloom Filter 交換記錄",
              "封包收發 byte 數",
              "80 條即時 debug log"])

    add_card(slide, 9.1, 4.1, 3.7, 2.5, "🔥 SOS + 火警",
             ["SOS_YELLOW / SOS_RED 廣播",
              "Triage Queue 頻寬搶佔 (QoS)",
              "FIRE_ALARM_RF 事件類型",
              "Si4467 433MHz RF 轉譯",
              "速率限制防訊號風暴"])

    # ── Slide 7: 程式碼結構 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_textbox(slide, 0.5, 0.3, 12, 0.8, "程式碼架構 (lib/ 目錄)", 32, ACCENT_ORANGE, True)

    table_data = [
        ["目錄/檔案", "行數", "職責"],
        ["main.dart", "~750", "App 入口、UI 整合、Provider 狀態管理"],
        ["mesh/ble_manager.dart", "774", "BLE 掃描+連線+同步 (Central 角色)"],
        ["mesh/event_manager.dart", "~700", "EventManager + TriageQueue 消費者"],
        ["mesh/mesh_event_handler.dart", "356", "接收解碼、去重、HLC merge、DB 寫入"],
        ["mesh/event_serializer.dart", "275", "Protobuf 序列化/反序列化工具"],
        ["mesh/mesh_router.dart", "113", "Zone-Based Geo-Fencing 路由決策"],
        ["mesh/native_bridge.dart", "~300", "MethodChannel ↔ Android Nordic BLE"],
        ["mesh/mesh_constants.dart", "53", "GATT UUID、BLE 連線參數常數"],
        ["crdt/hlc.dart", "89", "HLC 混合邏輯時鐘"],
        ["crdt/conflict_resolver.dart", "40", "Double-Spending 衝突解決"],
        ["crypto/identity_manager.dart", "124", "Ed25519 金鑰管理 + Trust Ladder"],
        ["crypto/signer.dart", "~50", "Payload 簽章工具"],
        ["proto/mesh_protocol.pb.dart", "~1200", "Protobuf 自動生成"],
        ["db/database_helper.dart", "~200", "SQLite 資料庫初始化"],
        ["geo/village_geofence.dart", "~300", "行政區邊界查詢"],
    ]
    add_table_slide(slide, 0.5, 1.3, 12.3, table_data, col_widths=[3.5, 1.0, 7.8], font_size=11)

    # ── Slide 8: 未來改進 (hive-btle 借鑑) ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_textbox(slide, 0.5, 0.3, 12, 0.8, "後續改進方向 — 借鑑 hive-btle 設計", 32, ACCENT_ORANGE, True)

    add_card(slide, 0.5, 1.3, 6.0, 2.5, "🥇 Power Profile 省電排程 (1~2 天)",
             ["現況: 掃 30s → 停 5s (duty ≈ 86%)",
              "改進: 依情境動態切換 scan interval",
              "  Aggressive: 掃 30s/停 5s (SOS 模式)",
              "  LowPower: 掃 5s/停 25s (背景待機)",
              "  UltraLow: 掃 3s/停 120s (低電量)",
              "預期: 電池從 ~3h → 12h+"])

    add_card(slide, 6.8, 1.3, 6.0, 2.5, "🥈 Advertising Metadata (2~3 天)",
             ["現況: 只廣播 Service UUID",
              "改進: 在 advertising data 中加入:",
              "  → Tier level (1 byte)",
              "  → 電量百分比 (1 byte)",
              "  → 有無待傳 SOS (1 bit)",
              "效果: 掃描階段即可篩選，減少無效連線"])

    add_card(slide, 0.5, 4.1, 6.0, 2.8, "🥉 Mesh-wide 加密 (3~5 天)",
             ["ChaCha20-Poly1305 AEAD 加密",
              "  → overhead 僅 30 bytes/封包",
              "Per-peer E2EE (X25519 key exchange)",
              "  → overhead 僅 46 bytes/封包",
              "  → 適合未來私人聊天功能",
              "兩層疊加 76 bytes，MTU 517 完全可行",
              "Dart cryptography 套件原生支援"])

    add_card(slide, 6.8, 4.1, 6.0, 2.8, "4️⃣ Leaf Node 被動模式 (1 天)",
             ["Tier 3 節點停止掃描",
              "只保留 GATT Server 被動接收",
              "由 Tier 1/2 主動連接並 Notify 推送",
              "",
              "5️⃣ Coded PHY 長距離 (硬體端)",
              "BLE 5.0 Coded PHY → 300m+ 範圍",
              "nRF54H20 基地台原生支援",
              "Flutter 端待 flutter_blue_plus 更新"])

    # ── Slide 9: 聊天室規劃 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_textbox(slide, 0.5, 0.3, 12, 0.8, "後續功能: 聊天室 — 基於現有架構擴充", 32, ACCENT_ORANGE, True)

    add_card(slide, 0.5, 1.3, 6.0, 2.8, "💬 設計方案",
             ["聊天訊息 = append-only event log",
              "  → 跟現有 Event_Logs 完全同構",
              "  → 不需要 Automerge CRDT",
              "",
              "新增 EventType: CHAT_MESSAGE = 8",
              "新增 ChatMessageData sub-message:",
              "  message_id, channel_id, content,",
              "  sent_at (HLC), display_name"])

    add_card(slide, 6.8, 1.3, 6.0, 2.8, "📐 封包估算",
             ["100 字中文 ≈ 300 bytes UTF-8",
              "+ MeshEvent header ≈ 150 bytes",
              "= 總計 ~450 bytes",
              "",
              "✅ MTU 517 內，無需分片",
              "✅ 走現有 Bloom Filter 同步",
              "✅ HLC 天然支援聊天排序",
              "✅ 無衝突 (append-only)"])

    add_card(slide, 0.5, 4.4, 12.3, 2.5, "🔧 實作路徑",
             ["1. proto: 新增 CHAT_MESSAGE EventType + ChatMessageData sub-message",
              "2. event_serializer.dart: 新增 buildChatEvent() 方法",
              "3. mesh_event_handler.dart: 新增 handleChatEvent() 處理邏輯",
              "4. UI: 新增聊天頁面 (channel list + message list + input)",
              "5. DB: Event_Logs 表原封不動，聊天訊息作為一種 event type 存入",
              "6. 如需私聊: 搭配 Per-peer E2EE (X25519 + ChaCha20) 加密",
              "預計工程量: 3~5 天 (不含 E2EE)"])

    # ── Slide 10: 結語 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_textbox(slide, 1, 1.5, 11, 1.2, "技術已驗證，等待硬體整合", 42, ACCENT_ORANGE, True, PP_ALIGN.CENTER)

    table_data = [
        ["模組", "狀態", "驗證方式"],
        ["BLE Mesh 多節點通訊", "✅ 完成", "實機跨廠牌測試 + debug log"],
        ["Bloom Filter 差異同步", "✅ 完成", "實機 log 分析"],
        ["HLC + CRDT 離線同步", "✅ 完成", "單元測試 + 整合測試"],
        ["Triage Queue QoS", "✅ 完成", "模擬測試"],
        ["地理圍欄路由", "✅ 完成", "程式邏輯驗證"],
        ["物資媒合 + 4-PIN 交割", "✅ 完成", "功能測試"],
        ["聊天室功能", "📋 規劃中", "基於現有架構擴充"],
        ["加密傳輸層", "📋 規劃中", "ChaCha20-Poly1305 + X25519"],
        ["省電排程", "📋 規劃中", "Power Profile 動態切換"],
    ]
    add_table_slide(slide, 2.0, 3.0, 9.3, table_data, col_widths=[3.0, 1.5, 4.8])

    output_path = os.path.join(OUTPUT_DIR, "烽傳_IgniRelay_App技術簡報.pptx")
    prs.save(output_path)
    print(f"[OK] PPTX #1 saved: {output_path}")
    return output_path


# ══════════════════════════════════════════════════════════════
#  PPTX #2: 整體系統框架 — 全局視角
# ══════════════════════════════════════════════════════════════
def generate_system_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── Slide 1: Cover ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_textbox(slide, 1, 0.5, 11, 1.2, "烽傳 IgniRelay", 48, ACCENT_ORANGE, True, PP_ALIGN.CENTER)
    add_textbox(slide, 1, 1.7, 11, 0.8, "去中心化災難通訊骨幹網路 — 系統整體架構", 28, WHITE, False, PP_ALIGN.CENTER)
    add_textbox(slide, 1, 3.0, 11, 0.5, "從中央到最基層里民的備援通訊生命線", 18, LIGHT_GRAY, False, PP_ALIGN.CENTER)
    add_textbox(slide, 1, 5.0, 11, 0.5, "LoRa × Wi-Fi HaLow × BLE × 433MHz RF — 四頻互補設計", 16, ACCENT_BLUE, False, PP_ALIGN.CENTER)
    add_textbox(slide, 1, 6.0, 11, 0.5, "Phase 0 — 概念驗證與 MVP  |  2026 年 3 月", 14, LIGHT_GRAY, False, PP_ALIGN.CENTER)

    # ── Slide 2: 問題定義 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_textbox(slide, 0.5, 0.3, 12, 0.8, "問題：縣市以下通訊斷層", 32, ACCENT_ORANGE, True)

    add_card(slide, 0.5, 1.3, 5.8, 2.5, "📊 現況數據",
             ["0403 花蓮強震: 172 座基地台受影響",
              "震後第 4 天仍有 46 座無法修復",
              "颱風後部分鄉鎮通訊孤島長達 16 天",
              "數位部僅一輛網路行動車可用"])

    add_card(slide, 6.6, 1.3, 6.2, 2.5, "🔍 核心斷層",
             ["政府防災通訊: 中央 ↔ 縣市 (衛星)",
              "✅ 已有方案",
              "",
              "縣市 ↔ 鄉鎮/消防/里長/避難所/民眾",
              "❌ 完全空白 — 烽傳填補此段"])

    add_card(slide, 0.5, 4.1, 12.3, 2.8, "🎯 系統目標",
             ["提供全台民眾與政府至少 14 天不中斷基礎通訊",
              "通訊從縣市 EOC 延伸至 → 鄉鎮公所 → 消防分隊 → 派出所 → 里長 → 避難所 → 一般民眾手機",
              "對齊台灣 EOC 六層行政指揮體系",
              "不依賴外國通訊基礎設施運作 (離網自主)",
              "零門檻接入: 民眾只需手機安裝 App"])

    # ── Slide 3: 四頻互補設計 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_textbox(slide, 0.5, 0.3, 12, 0.8, "頻段互補策略：讓不同頻段做最擅長的事", 32, ACCENT_ORANGE, True)

    table_data = [
        ["頻段", "頻率", "系統角色", "關鍵規格"],
        ["LoRa (LR2021)", "920 MHz", "控制面：SOS 直達、喚醒 HaLow、心跳", "5-50km、-141.5 dBm、低功耗"],
        ["Wi-Fi HaLow (MM8108)", "920-928 MHz", "資料面：視訊通話、圖片傳輸、指揮同步", "4MHz ~21 Mbps、TWT 0.76mA"],
        ["BLE 6.0 (nRF54H20)", "2.4 GHz", "民眾終端接入、手機 App、節點管理", "100m、Channel Sounding 定位"],
        ["433 MHz RF (Si4467)", "433 MHz", "被動接收住警器火警訊號", "穿透鋼筋混凝土、零汰換"],
    ]
    add_table_slide(slide, 0.5, 1.4, 12.3, table_data, col_widths=[2.5, 1.5, 4.5, 3.8], font_size=12)

    add_card(slide, 0.5, 4.2, 6.0, 2.8, "🔀 Transport 分流邏輯",
             ["SOS_RED → LoRa 優先 (佇列>10 升級 HaLow)",
              "SOS_YELLOW / HAZARD → LoRa 優先",
              "RESOURCE / INFO → HaLow 批次 (30s 累積)",
              "下行廣播 → 小封包 LoRa / 大檔 LoRa+HaLow",
              "民眾 ↔ 節點 → BLE only",
              "視訊通話 → HaLow 專用"])

    add_card(slide, 6.8, 4.2, 6.0, 2.8, "📡 HaLow 喚醒機制",
             ["TWT 待機: 平均 0.76 mA",
              "LoRa 喚醒封包攜帶路由資訊:",
              "  → GPS 座標 + 所需頻寬",
              "  → 節點判斷是否為最佳中繼",
              "TWT 窗口: 200ms ~ 5s",
              "HaLow association: ~80ms"])

    # ── Slide 4: 硬體節點 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_textbox(slide, 0.5, 0.3, 12, 0.8, "四種硬體節點 — 統一 nRF54H20 主控", 32, ACCENT_ORANGE, True)

    table_data = [
        ["節點類型", "通訊模組", "部署位置", "供電", "估算成本"],
        ["節點一 (道路中繼)", "nRF54H20+LR2021+MM8108×2", "道路沿線桿掛 (每 300-500m)", "10W太陽能+LFP 40Ah", "NT$15,500-25,500"],
        ["節點二 (政府機關)", "同節點一 (+Si4467選配)", "公所/消防/派出所/醫院", "市電UPS+LFP 100Ah", "NT$13,000-22,000"],
        ["節點三 (室內擴展)", "nRF54H20+nRF7002+Ethernet", "避難所/里長室 室內", "市電UPS+LFP 20Ah", "NT$6,000-10,000"],
        ["節點四 (末端感知)", "nRF54H20+Si4467 433MHz", "全台路燈/電線桿", "市電UPS+LFP 5Ah", "NT$2,500-4,000"],
    ]
    add_table_slide(slide, 0.5, 1.3, 12.3, table_data, col_widths=[2.0, 3.5, 3.0, 2.0, 1.8], font_size=11)

    add_card(slide, 0.5, 4.3, 6.0, 2.6, "🔋 14 天自持設計",
             ["全系統 LFP (磷酸鐵鋰) + Buck 降壓 >90%",
              "節點一: 太陽能補充 560Wh/14天 → 餘量 >2.3×",
              "節點二: 1,200Wh 可用 → 餘量 1.5×",
              "節點三: 240Wh 可用 → 餘量 1.8×",
              "節點四: 64Wh 可用 → 餘量 1.9×"])

    add_card(slide, 6.8, 4.3, 6.0, 2.6, "🏗️ 部署策略",
             ["沿路網鋪設: 省道/縣道/產業道路",
              "每 300-500m 一個節點 (HaLow 有效距離)",
              "覆蓋沿線公所、警消、里與避難所",
              "不需高山鐵塔施工",
              "HaLow 雙模組: 上行 21Mbps / 下行 21Mbps"])

    # ── Slide 5: 六層指揮通訊 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_textbox(slide, 0.5, 0.3, 12, 0.8, "六層指揮通訊與視訊頻寬規劃", 32, ACCENT_ORANGE, True)

    table_data = [
        ["通訊層", "連線對象", "頻寬需求", "技術方案"],
        ["縣市 EOC ↔ 鄉鎮公所", "同時 1-2 路 720p 發言", "~10-14 Mbps", "MM8108×2 (各21 Mbps)"],
        ["縣市 ↔ 368 鄉鎮公所", "主動發言才上傳 720p", "~2-4 Mbps", "MM8108×1 (4MHz)"],
        ["鄉鎮 ↔ 里/避難所", "單向視訊 + 雙向語音", "~1-2 Mbps", "HaLow MCS2-7"],
        ["縣市 ↔ 警消 (網狀)", "雙向通訊", "~4-8 Mbps", "MM8108×2 網狀冗餘"],
        ["分局/大隊 ↔ 基層", "語音通話 + 圖片", "0.3-1 Mbps", "HaLow MCS0-2"],
        ["里/避難所 ↔ 民眾", "文字、SOS、物資媒合", "BLE 頻寬", "BLE Mesh + Bloom Filter"],
    ]
    add_table_slide(slide, 0.5, 1.3, 12.3, table_data, col_widths=[2.5, 3.0, 2.0, 4.8], font_size=12)

    add_card(slide, 0.5, 4.8, 6.0, 2.2, "📹 動態視訊品質",
             ["發言中: 720p H.265 ~720 kbps",
              "待機聆聽: 360p H.265 ~0.5 Mbps",
              "純語音: 16-64 kbps (Opus 編碼)",
              "頻寬需求: 32 Mbps → 實際 10-14 Mbps"])

    add_card(slide, 6.8, 4.8, 6.0, 2.2, "⏱️ 延遲估算",
             ["HaLow 單跳: 5-15 ms",
              "10km 多跳: ~100ms",
              "50km 多跳: ~500ms (≈ GEO 衛星)",
              "HaLert 論文實測: 15-54.8ms"])

    # ── Slide 6: 核心演算法 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_textbox(slide, 0.5, 0.3, 12, 0.8, "七大核心通訊演算法", 32, ACCENT_ORANGE, True)

    algos = [
        ("⏰ HLC 混合邏輯時鐘", "防斷電歸零、因果排序不錯亂"),
        ("🔄 Bloom Filter 差異同步", "5,000容量/1%FPR/5.85KB，僅傳差量"),
        ("🤝 CRDT 離線同步", "孤島各自修改，重逢時最終一致"),
        ("🚨 Triage Queue 搶佔", "SOS_RED 100% 頻寬搶佔"),
        ("🗺️ Zone-Based Geo-Fencing", "里→鄉鎮市區分層路由"),
        ("🛡️ Quarantine Vote", "加權投票去中心化隔離惡意節點"),
        ("🔥 433 MHz RF 轉譯", "住警器零汰換智慧升級"),
    ]
    for i, (title, desc) in enumerate(algos):
        col = i % 3
        row = i // 3
        add_card(slide, 0.5 + col * 4.2, 1.3 + row * 2.8, 3.9, 2.5,
                 title, [desc, "", "✅ Phase 0 已實作驗證" if i < 6 else "Phase 1 硬體整合"])

    # ── Slide 7: 管理操作介面 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_textbox(slide, 0.5, 0.3, 12, 0.8, "政府管理介面功能規劃 (Dashboard)", 32, ACCENT_ORANGE, True)

    add_card(slide, 0.5, 1.3, 4.0, 2.8, "📊 即時監控",
             ["全台節點健康度地圖",
              "各節點電量/溫度/連線數",
              "鄉鎮層級 SOS 統計摘要",
              "LoRa 骨幹 RSSI 品質圖",
              "HaLow 吞吐量即時監測"])

    add_card(slide, 4.8, 1.3, 4.0, 2.8, "🚨 災情管理",
             ["SOS 事件列表 + 地圖標記",
              "危險區域圖層 (Hazard overlay)",
              "物資供需統計報表",
              "火警 RF 觸發即時通報",
              "Quarantine 惡意節點追蹤"])

    add_card(slide, 9.1, 1.3, 3.7, 2.8, "👤 身份管理",
             ["信任等級審核 (L1-L3)",
              "TW FidO 實名認證介面",
              "Quarantine 投票覆核",
              "黑名單管理",
              "社群背書 (L2 驗證)"])

    add_card(slide, 0.5, 4.4, 4.0, 2.5, "⚙️ 系統設定",
             ["節點 OTA 韌體更新",
              "Bloom Filter 參數調整",
              "Triage Queue 閾值設定",
              "Zone 路由邊界編輯",
              "加密金鑰輪換管理"])

    add_card(slide, 4.8, 4.4, 4.0, 2.5, "📈 統計分析",
             ["歷史事件查詢 + 匯出",
              "各里/鄉鎮通訊覆蓋率報告",
              "物資媒合成功率分析",
              "節點電量趨勢預測",
              "災情演習報告生成"])

    add_card(slide, 9.1, 4.4, 3.7, 2.5, "📡 通訊管理",
             ["HaLow 頻道分配管理",
              "LoRa 骨幹路徑最佳化",
              "視訊通話權限控制",
              "頻寬配額管理",
              "下行廣播發布"])

    # ── Slide 8: 全台部署規模 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_textbox(slide, 0.5, 0.3, 12, 0.8, "全台部署規模估算", 32, ACCENT_ORANGE, True)

    table_data = [
        ["關鍵設施 (錨點)", "全台數量", "資料來源"],
        ["鄉鎮市區公所", "368 處", "內政部 114 年 12 月底"],
        ["消防分隊 (含港務)", "628 隊", "消防署 112 年統計"],
        ["警察分駐所及派出所", "1,491 所", "警政署 2024 年 7 月"],
        ["縣市 EOC/警察局/消防局", "66 處", "行政院防災體系 (22×3)"],
        ["災害避難收容所", "5,928 處", "內政部 2025 年 6 月"],
        ["重大醫院", "~102 家", "衛福部 2025 年評鑑"],
        ["里長/村長辦公室", "7,752 處", "內政部 114 年 12 月底"],
        ["錨點合計", "~15,935 個", ""],
    ]
    add_table_slide(slide, 0.5, 1.3, 8.0, table_data, col_widths=[2.5, 1.5, 4.0], font_size=12)

    add_card(slide, 8.8, 1.3, 4.0, 5.2, "📅 分階段建置",
             ["Phase 0 (現況)",
              "  BLE Mesh App MVP",
              "  成本: NT$0",
              "",
              "Phase 1 (鄉鎮試點)",
              "  單一偏遠鄉鎮 ~35 錨點",
              "  約 NT 50-80 萬 (硬體)",
              "",
              "Phase 2 (縣市試點)",
              "  單縣市 ~650 錨點",
              "  約 NT 800-1,300 萬 (硬體)",
              "",
              "Phase 3 (全台願景)",
              "  22 縣市 ~8,000 關鍵設施",
              "  待實測後修正"])

    # ── Slide 9: 住警器 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_textbox(slide, 0.5, 0.3, 12, 0.8, "創新亮點：住警器零汰換智慧升級", 32, ACCENT_ORANGE, True)

    add_card(slide, 0.5, 1.3, 6.0, 2.5, "🔥 433 MHz RF 住警器轉譯",
             ["Si4467 被動接收住警器 RF 訊號",
              "支援: 宏力、TYY 等台灣主流品牌",
              "編碼格式: EV1527 / PT2262",
              "偵測 → 轉譯 → FIRE_ALARM_RF 事件",
              "自動標記 SOS_RED + 基地台 GPS"])

    add_card(slide, 6.8, 1.3, 6.0, 2.5, "💡 民眾效益",
             ["零成本: 不需汰換現有住警器",
              "433 MHz 穿透鋼筋混凝土 > BLE",
              "直通消防局的智慧物聯網升級",
              "火勢蔓延前精確位置預警",
              "速率限制防地震後集體誤觸發"])

    # ── Slide 10: 結語 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_textbox(slide, 1, 1.0, 11, 1.2, "烽傳 IgniRelay", 48, ACCENT_ORANGE, True, PP_ALIGN.CENTER)
    add_textbox(slide, 1, 2.2, 11, 0.8, "Phase 0 技術驗證完成，等待 Phase 1 硬體整合", 24, WHITE, False, PP_ALIGN.CENTER)

    add_bullet_list(slide, 2.0, 3.5, 9.0, 2.5, [
        "BLE Mesh / HLC / CRDT / Bloom Filter / Triage Queue / Geo-Fencing — 全數實作驗證完成",
        "HaLert 論文 (2025) 以同架構實測驗證可行性: 延遲 15-54.8ms、吞吐量 134-726 Kbps",
        "Phase 1 NT$50-80 萬即可取得台灣地形下三頻骨幹實測數據",
        "對齊行政院強韌台灣計畫，填補縣市以下通訊真空",
    ], font_size=16, color=LIGHT_GRAY)

    add_textbox(slide, 1, 6.0, 11, 0.8,
                "「烽火台是古代第一套 Mesh 網路。烽傳，就是把這個概念裝進每個人的口袋。」",
                18, ACCENT_BLUE, False, PP_ALIGN.CENTER)

    output_path = os.path.join(OUTPUT_DIR, "烽傳_IgniRelay_系統整體架構簡報.pptx")
    prs.save(output_path)
    print(f"[OK] PPTX #2 saved: {output_path}")
    return output_path


if __name__ == "__main__":
    p1 = generate_app_pptx()
    p2 = generate_system_pptx()
    print(f"\nDone! Files saved to:\n  {p1}\n  {p2}")
